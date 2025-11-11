import chainlit as cl
import pandas as pd
import docx
from pptx import Presentation
import os
import json
from bs4 import BeautifulSoup
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_community.document_loaders import PyMuPDFLoader, BSHTMLLoader

# Disable telemetry to avoid config issues
os.environ["CHAINLIT_TELEMETRY_ENABLED"] = "false"

# Initialize the LLM using Hugging Face secrets
try:
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash",
        google_api_key=os.environ.get('GOOGLE_API_KEY'), 
        temperature=0.1
    )
except Exception as e:
    llm = None
    print(f"LLM initialization error: {e}")

@cl.on_chat_start
async def start():
    if not llm:
        await cl.Message(content="Error: `GOOGLE_API_KEY` is not properly configured. Please check the space settings.").send()
        return

    await cl.Message(content="Hello! Please upload a file to begin (CSV, XLSX, DOCX, PDF, PPTX, TXT, HTML, XML, JSON).").send()
    
    files = await cl.AskFileMessage(
        content="Please upload a file to ask questions about its content.",
        accept=[
            "text/csv",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", # XLSX
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document", # DOCX
            "application/pdf", # PDF
            "application/vnd.openxmlformats-officedocument.presentationml.presentation", # PPTX
            "text/plain", # TXT
            "text/html", # HTML
            "text/xml", "application/xml", # XML
            "application/json", # JSON
        ],
        max_size_mb=10,  # Reduced for Hugging Face
        timeout=180,
    ).send()
    
    if not files:
        await cl.Message(content="No file uploaded. Please restart the chat to try again.").send()
        return

    file = files[0]
    msg = cl.Message(content=f"Processing `{file.name}`...")
    await msg.send()

    try:
        # Process different file types
        if file.type == "text/csv" or file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            if file.type == "text/csv":
                df = pd.read_csv(file.path)
            else:
                df = pd.read_excel(file.path)
            
            cl.user_session.set("data", df)
            cl.user_session.set("file_type", "csv")
            msg.content = f"`{file.name}` uploaded successfully. You can now ask questions about your data."
            await msg.update()
        
        else:
            full_text = ""
            if file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                document = docx.Document(file.path)
                full_text = "\n".join([para.text for para in document.paragraphs])

            elif file.type == "application/pdf":
                loader = PyMuPDFLoader(file.path)
                documents = loader.load()
                full_text = "\n".join([doc.page_content for doc in documents])

            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(file.path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                full_text = "\n".join(text_runs)
            
            elif file.type == "text/plain":
                with open(file.path, "r", encoding="utf-8") as f:
                    full_text = f.read()
            
            elif file.type == "text/html":
                loader = BSHTMLLoader(file.path)
                documents = loader.load()
                full_text = "\n".join([doc.page_content for doc in documents])

            elif file.type in ["text/xml", "application/xml"]:
                with open(file.path, "r", encoding="utf-8") as f:
                    content = f.read()
                soup = BeautifulSoup(content, "xml")
                full_text = soup.get_text(separator="\n", strip=True)

            elif file.type == "application/json":
                with open(file.path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                full_text = json.dumps(data, indent=2)

            cl.user_session.set("data", full_text)
            cl.user_session.set("file_type", "text")
            msg.content = f"`{file.name}` uploaded successfully. You can now ask questions."
            await msg.update()

    except Exception as e:
        msg.content = f"Error processing file: {str(e)}"
        await msg.update()

@cl.on_message
async def main(message: cl.Message):
    data = cl.user_session.get("data")
    file_type = cl.user_session.get("file_type")
    file_name = cl.user_session.get("file_name", "Uploaded File")
    
    if data is None or not llm:
        await cl.Message(content="The application is not ready. Please ensure a file has been uploaded.").send()
        return

    thinking_msg = cl.Message(content="Thinking...")
    await thinking_msg.send()

    try:
        if file_type == "csv":
            # For CSV files, provide dataframe info to the LLM
            if isinstance(data, pd.DataFrame):
                shape = data.shape
                columns = list(data.columns)
                head = data.head(3).to_string()
                
                prompt_template = """
                Analyze this dataset and answer the question.
                
                Dataset Info:
                - Rows: {rows}, Columns: {cols}
                - Columns: {columns}
                - Sample data:
                {sample}
                
                Question: {question}
                
                Provide a clear, helpful response based on the data.
                """
                
                prompt = ChatPromptTemplate.from_template(prompt_template)
                chain = prompt | llm | StrOutputParser()
                response = await cl.make_async(chain.invoke)({
                    "rows": shape[0],
                    "cols": shape[1],
                    "columns": columns,
                    "sample": head,
                    "question": message.content
                })
            else:
                response = "Error: Data format is not recognized as a DataFrame."
                
        else:  # text files
            # Limit context length to avoid token limits
            context = str(data)[:4000] if len(str(data)) > 4000 else str(data)
            
            prompt_template = """
            Analyze this document content and answer the question.
            
            Content:
            {context}
            
            Question: {question}
            
            Provide a helpful response based on the content above.
            If you cannot answer from the content, please state that clearly.
            """
            
            prompt = ChatPromptTemplate.from_template(prompt_template)
            chain = prompt | llm | StrOutputParser()
            response = await cl.make_async(chain.invoke)({
                "context": context,
                "question": message.content
            })
        
        thinking_msg.content = response
        await thinking_msg.update()
            
    except Exception as e:
        error_message = f"An error occurred: {str(e)}"
        thinking_msg.content = error_message
        await thinking_msg.update()